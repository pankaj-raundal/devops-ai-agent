"""Generate a professional presentation for DevOps AI Agent."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (Professional Light Theme) ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x2D, 0x3A, 0x4A)
SUBTITLE_TEXT = RGBColor(0x5A, 0x6A, 0x7E)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
ACCENT_GREEN = RGBColor(0x0D, 0x94, 0x53)
ACCENT_ORANGE = RGBColor(0xE8, 0x71, 0x0A)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
ACCENT_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)
BORDER_GRAY = RGBColor(0xDA, 0xDE, 0xE3)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_box(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=DARK_TEXT, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
        # Bullet character
        p.bullet = True  # type: ignore[assignment]
    return txBox


def add_icon_card(slide, left, top, width, height, icon, title, desc,
                  bg_color=WHITE, accent_color=ACCENT_BLUE):
    """Add a card with an icon circle, title, and description."""
    box = add_shape_box(slide, left, top, width, height, bg_color, BORDER_GRAY)

    # Icon circle
    circle_size = Inches(0.55)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, left + Inches(0.15), top + Inches(0.85), width - Inches(0.3), Inches(0.35),
                 title, font_size=12, bold=True, color=DARK_TEXT)
    # Desc
    add_text_box(slide, left + Inches(0.15), top + Inches(1.15), width - Inches(0.3), height - Inches(1.35),
                 desc, font_size=10, color=SUBTITLE_TEXT)
    return box


def slide_header(slide, title, subtitle=None):
    """Add a consistent header bar to content slides."""
    # Blue accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, Inches(0.7), Inches(0.25), Inches(10), Inches(0.55),
                 title, font_size=28, bold=True, color=DARK_TEXT)
    if subtitle:
        add_text_box(slide, Inches(0.7), Inches(0.75), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=SUBTITLE_TEXT)


# ── SLIDES ──

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Accent band at top
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT_BLUE
    band.line.fill.background()

    # Main title
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
                 "DevOps AI Agent", font_size=44, bold=True, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
                 "Automated Story-to-Code Pipeline", font_size=22,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Thin divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.55), Inches(2.33), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = BORDER_GRAY
    div.line.fill.background()

    # Description
    add_text_box(slide, Inches(2), Inches(3.85), Inches(9.33), Inches(0.8),
                 "AI-powered automation that reads Azure DevOps stories, implements code changes,\n"
                 "runs tests, and performs code review — reducing developer cycle time by up to 60%.",
                 font_size=14, color=SUBTITLE_TEXT, alignment=PP_ALIGN.CENTER)

    # Bottom info
    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
                 "Presented by: Pankaj Raundal  |  Team: VTD Producer Connectors  |  March 2026",
                 font_size=12, color=SUBTITLE_TEXT, alignment=PP_ALIGN.CENTER)

    # Bottom accent band
    band2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12))
    band2.fill.solid()
    band2.fill.fore_color.rgb = ACCENT_BLUE
    band2.line.fill.background()


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "Agenda")

    items = [
        ("1", "Problem Statement", "Current challenges in the development workflow"),
        ("2", "Solution Overview", "What DevOps AI Agent does and how it works"),
        ("3", "Architecture & Pipeline", "Technical design, dual AI strategy, and flow"),
        ("4", "Key Features", "Plan-Approve-Execute, smart merge, queue, and more"),
        ("5", "Live Dashboard", "Real-time monitoring, plan approval, and control"),
        ("6", "Benefits & ROI", "Measurable impact on team productivity"),
        ("7", "Roadmap & Next Steps", "Path to production readiness"),
    ]

    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.4) + Inches(i * 0.72)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(1.85), y - Inches(0.02), Inches(4), Inches(0.35),
                     title, font_size=16, bold=True, color=DARK_TEXT)
        add_text_box(slide, Inches(1.85), y + Inches(0.28), Inches(8), Inches(0.3),
                     desc, font_size=12, color=SUBTITLE_TEXT)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "The Problem", "Current developer workflow is manual and repetitive")

    pain_points = [
        ("⏱", "Time-Consuming", "Developers spend 30-45 min per story just on\nsetup: read story, create branch, scaffold code.", ACCENT_ORANGE),
        ("🔁", "Repetitive Tasks", "Same pattern every time: fetch story → branch →\nimplement → test → review → commit.", ACCENT_BLUE),
        ("🐛", "Human Error", "Manual processes lead to missed acceptance criteria,\ninconsistent branch naming, skipped tests.", ACCENT_PURPLE),
        ("📊", "No Visibility", "Managers lack insight into implementation progress.\nNo real-time tracking or audit trail.", ACCENT_TEAL),
    ]

    for i, (icon, title, desc, color) in enumerate(pain_points):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + Inches(col * 5.8)
        top = Inches(1.6) + Inches(row * 2.3)
        add_icon_card(slide, left, top, Inches(5.3), Inches(1.9), icon, title, desc, WHITE, color)

    # Bottom callout
    callout = add_shape_box(slide, Inches(0.8), Inches(6.2), Inches(11.73), Inches(0.7),
                            LIGHT_BLUE_BG, ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.5),
                 "Result: Slower delivery, inconsistent quality, and developer frustration.",
                 font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "The Solution: DevOps AI Agent",
                 "End-to-end automation from story assignment to review-ready code")

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                 "A Python CLI tool + web dashboard that automates the entire developer workflow:",
                 font_size=15, color=DARK_TEXT)

    steps = [
        ("📋", "Read Story", "Fetches assigned story\nfrom Azure DevOps\n(or Zendesk ticket)"),
        ("🔬", "AI Analysis", "Analyzes complexity,\nrisk, and determines\nif code changes needed"),
        ("🌿", "Create Branch", "Auto-creates properly\nnamed feature branch\nfollowing conventions"),
        ("📝", "Plan & Approve", "AI generates a plan;\nhuman reviews per-file\nbefore any code is written"),
        ("✅", "Test & Review", "Runs tests (PHPUnit,\nPHPCS, PHPStan) +\nAI code review"),
    ]

    for i, (icon, title, desc) in enumerate(steps):
        left = Inches(0.5) + Inches(i * 2.45)
        top = Inches(2.35)
        w = Inches(2.15)
        h = Inches(2.4)

        add_icon_card(slide, left, top, w, h, icon, title, desc, WHITE, ACCENT_BLUE)

        # Arrow between cards
        if i < len(steps) - 1:
            arrow_x = left + w + Inches(0.05)
            add_text_box(slide, arrow_x, top + Inches(0.8), Inches(0.3), Inches(0.4),
                         "→", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Key stat boxes
    stats = [
        ("60%", "Faster Delivery", ACCENT_GREEN),
        ("100%", "Test Coverage", ACCENT_BLUE),
        ("24/7", "Available", ACCENT_TEAL),
    ]
    for i, (val, label, color) in enumerate(stats):
        left = Inches(2.0) + Inches(i * 3.5)
        top = Inches(5.4)
        box = add_shape_box(slide, left, top, Inches(2.5), Inches(1.1), WHITE, color)
        add_text_box(slide, left + Inches(0.1), top + Inches(0.1), Inches(2.3), Inches(0.55),
                     val, font_size=28, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), top + Inches(0.6), Inches(2.3), Inches(0.35),
                     label, font_size=13, color=SUBTITLE_TEXT, alignment=PP_ALIGN.CENTER)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "Architecture & Pipeline Flow",
                 "Modular, extensible Python architecture")

    # Pipeline flow - horizontal
    flow_steps = [
        ("Zendesk\nTicket", ACCENT_ORANGE),
        ("Azure DevOps\nStory", ACCENT_BLUE),
        ("AI Analysis\n& Triage", ACCENT_PURPLE),
        ("Plan &\nApprove", ACCENT_TEAL),
        ("Implement\n(CLI/API)", ACCENT_GREEN),
        ("Test &\nReview", ACCENT_BLUE),
        ("Feature\nBranch ✓", ACCENT_GREEN),
    ]

    for i, (label, color) in enumerate(flow_steps):
        left = Inches(0.3) + Inches(i * 1.82)
        top = Inches(1.5)
        w = Inches(1.5)
        h = Inches(0.9)
        box = add_shape_box(slide, left, top, w, h, color, None)
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"

        if i < len(flow_steps) - 1:
            add_text_box(slide, left + w - Inches(0.05), top + Inches(0.2), Inches(0.4), Inches(0.4),
                         "▸", font_size=18, color=BORDER_GRAY)

    # Module structure
    add_text_box(slide, Inches(0.7), Inches(2.8), Inches(5), Inches(0.35),
                 "Module Structure", font_size=16, bold=True, color=DARK_TEXT)

    modules = [
        ("src/cli.py", "Click CLI — 8 commands (fetch, run, run-all, watch, webhook, dashboard...)"),
        ("src/pipeline.py", "End-to-end orchestrator with queue, plan-approve-execute flow"),
        ("src/agent/", "Analyzer, plan engine (append/replace merge), implementation agent"),
        ("src/integrations/", "Azure DevOps, Zendesk, Git (workspace reset), Webhook connectors"),
        ("src/reviewer/", "Test runner (PHPUnit/PHPCS/PHPStan) + AI code reviewer"),
        ("src/dashboard/", "Flask web UI with SSE streaming + plan approval modal"),
    ]

    for i, (mod, desc) in enumerate(modules):
        y = Inches(3.2) + Inches(i * 0.55)
        add_text_box(slide, Inches(0.9), y, Inches(2), Inches(0.35),
                     mod, font_size=12, bold=True, color=ACCENT_BLUE)
        add_text_box(slide, Inches(3.2), y, Inches(6), Inches(0.35),
                     desc, font_size=11, color=SUBTITLE_TEXT)

    # Tech stack on right
    add_text_box(slide, Inches(7.5), Inches(2.8), Inches(4), Inches(0.35),
                 "Technology Stack", font_size=16, bold=True, color=DARK_TEXT)

    stack_items = [
        ("Language", "Python 3.10+"),
        ("AI Provider", "GitHub Models API (GPT-4o)"),
        ("AI CLI", "Claude Code CLI / Codex CLI"),
        ("CLI", "Click + Rich"),
        ("Dashboard", "Flask + SSE"),
        ("DevOps", "Azure DevOps (az CLI)"),
        ("VCS", "GitPython"),
    ]

    for i, (key, val) in enumerate(stack_items):
        y = Inches(3.2) + Inches(i * 0.5)
        add_text_box(slide, Inches(7.7), y, Inches(1.8), Inches(0.3),
                     key, font_size=11, bold=True, color=DARK_TEXT)
        add_text_box(slide, Inches(9.6), y, Inches(2.5), Inches(0.3),
                     val, font_size=11, color=SUBTITLE_TEXT)


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "Key Features")

    features = [
        ("�", "Plan-Approve-Execute", "AI generates structured plan with\nper-file changes. Human reviews and\napproves before any code is written.", ACCENT_PURPLE),
        ("🔀", "Smart Merge Strategy", "CLI tools read files & replace.\nAPI returns only new code — pipeline\nappends to existing files safely.", ACCENT_BLUE),
        ("🔄", "State Transitions", "Automatically moves stories to\nTesting/Evaluation states to prevent\nre-processing.", ACCENT_GREEN),
        ("📋", "Story Queue", "Batch-process multiple stories\nwith workspace reset between each.\nFull per-story tracking.", ACCENT_TEAL),
        ("🛡", "Data Consent", "Security-first: explicit consent before\nsending any code or story data to\nexternal AI providers.", ACCENT_ORANGE),
        ("📊", "Write-back & Alerts", "Posts AI results back to Azure\nDevOps and Zendesk. Real-time\ntoast notifications + Slack.", ACCENT_BLUE),
        ("🔬", "AI Analysis & Triage", "Analyzes complexity, risk, and\naffected areas. Determines if code\nchanges are actually needed.", ACCENT_PURPLE),
        ("🏃", "Watch & Webhook", "Polls for new stories or receives\npush events. Auto-runs pipeline.\nSet it and forget it.", ACCENT_GREEN),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 4
        row = i // 4
        left = Inches(0.4) + Inches(col * 3.1)
        top = Inches(1.4) + Inches(row * 2.7)
        add_icon_card(slide, left, top, Inches(2.8), Inches(2.3), icon, title, desc, WHITE, color)


def slide_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "Live Dashboard", "Real-time monitoring and control via web interface")

    # Dashboard feature cards
    dash_features = [
        ("Real-Time Progress", "Server-Sent Events (SSE) stream\npipeline stages live to the browser.\n10-stage visual progress bar."),
        ("Plan Approval Modal", "Review AI-generated plans per-file.\nApprove, reject, or inspect code\nbefore any changes are applied."),
        ("Story Queue Panel", "View all queued stories, track status\n(queued/in-progress/done/failed).\nClick any story to see its details."),
        ("Pipeline Controls", "Fetch Story, Run Pipeline, Run All,\nDashboard — all with one click.\nToast alerts for every event."),
    ]

    for i, (title, desc) in enumerate(dash_features):
        left = Inches(0.5) + Inches(i * 3.1)
        top = Inches(1.5)
        box = add_shape_box(slide, left, top, Inches(2.8), Inches(2.0), LIGHT_BLUE_BG, ACCENT_BLUE)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.15), Inches(2.5), Inches(0.35),
                     title, font_size=14, bold=True, color=ACCENT_BLUE)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.55), Inches(2.5), Inches(1.3),
                     desc, font_size=11, color=DARK_TEXT)

    # CLI commands table
    add_text_box(slide, Inches(0.7), Inches(4.0), Inches(5), Inches(0.35),
                 "CLI Commands", font_size=16, bold=True, color=DARK_TEXT)

    commands = [
        ("dai fetch", "Fetch latest assigned story from Azure DevOps"),
        ("dai run [-s ID]", "Full pipeline: fetch → branch → implement → test → review"),
        ("dai run-all", "Queue all matching stories and process sequentially"),
        ("dai implement", "Implement current story (reads .current-story.md)"),
        ("dai review", "Run tests + AI code review on current changes"),
        ("dai watch", "Poll for new stories and auto-run pipeline"),
        ("dai dashboard", "Launch the web dashboard on port 8090"),
        ("dai webhook", "Start Flask webhook server for push events"),
    ]

    for i, (cmd, desc) in enumerate(commands):
        y = Inches(4.45) + Inches(i * 0.38)
        bg = LIGHT_BG if i % 2 == 0 else WHITE
        row_bg = add_shape_box(slide, Inches(0.7), y, Inches(11.5), Inches(0.36), bg, None)
        add_text_box(slide, Inches(0.9), y + Inches(0.03), Inches(2.5), Inches(0.3),
                     cmd, font_size=11, bold=True, color=ACCENT_BLUE)
        add_text_box(slide, Inches(3.8), y + Inches(0.03), Inches(8), Inches(0.3),
                     desc, font_size=11, color=DARK_TEXT)


def slide_benefits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "Benefits & ROI", "Measurable impact on team productivity and quality")

    # Before vs After comparison
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
                 "Before (Manual)", font_size=18, bold=True, color=ACCENT_ORANGE)
    add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5), Inches(0.4),
                 "After (AI Agent)", font_size=18, bold=True, color=ACCENT_GREEN)

    before_items = [
        "30-45 min per story for setup and scaffolding",
        "Manual branch creation and naming",
        "Developers read and interpret stories manually",
        "Tests often skipped under time pressure",
        "No visibility into implementation progress",
        "Inconsistent code review quality",
    ]
    after_items = [
        "5-10 min — AI handles setup and implementation",
        "Auto-created branches following naming conventions",
        "AI generates plan, human approves per-file before apply",
        "Automated test suite runs on every implementation",
        "Real-time dashboard with plan approval + story tracking",
        "CLI reads files directly; API appends safely (no wipes)",
    ]

    for i, (before, after) in enumerate(zip(before_items, after_items)):
        y = Inches(1.95) + Inches(i * 0.6)
        # Before
        add_shape_box(slide, Inches(0.7), y, Inches(5.5), Inches(0.5), WHITE, ACCENT_ORANGE)
        add_text_box(slide, Inches(0.9), y + Inches(0.06), Inches(5.1), Inches(0.38),
                     "✗  " + before, font_size=11, color=DARK_TEXT)
        # After
        add_shape_box(slide, Inches(6.9), y, Inches(5.5), Inches(0.5), WHITE, ACCENT_GREEN)
        add_text_box(slide, Inches(7.1), y + Inches(0.06), Inches(5.1), Inches(0.38),
                     "✓  " + after, font_size=11, color=DARK_TEXT)

    # ROI summary
    roi_box = add_shape_box(slide, Inches(0.7), Inches(5.8), Inches(11.73), Inches(1.2),
                            LIGHT_BLUE_BG, ACCENT_BLUE)
    roi_items = [
        ("~60%", "reduction in story cycle time"),
        ("~80%", "fewer manual repetitive tasks"),
        ("100%", "automated test + review coverage"),
        ("24/7", "pipeline availability"),
    ]
    for i, (val, desc) in enumerate(roi_items):
        left = Inches(1.0) + Inches(i * 3.0)
        add_text_box(slide, left, Inches(5.95), Inches(1.2), Inches(0.4),
                     val, font_size=24, bold=True, color=ACCENT_BLUE)
        add_text_box(slide, left + Inches(1.3), Inches(6.0), Inches(1.6), Inches(0.35),
                     desc, font_size=12, color=DARK_TEXT)


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "How It Works — Demo Flow")

    steps = [
        ("1", "Configure", "Set up Azure DevOps org, project,\nAI provider, and target project path\nin config.yaml", ACCENT_BLUE),
        ("2", "Assign Story", "Assign a story in Azure DevOps\nwith the 'auto' tag to trigger\nAI processing", ACCENT_PURPLE),
        ("3", "Run Pipeline", "Execute `dai run` or `dai run-all`\nor use the web dashboard\nto trigger processing", ACCENT_GREEN),
        ("4", "AI Plans", "AI reads story, generates structured\nplan with per-file changes, risk\nassessment, and testing steps", ACCENT_TEAL),
        ("5", "Approve & Apply", "Human reviews plan, approves or\nrejects per-file. CLI reads files\ndirectly; API appends safely.", ACCENT_ORANGE),
        ("6", "Test & Ship", "Automated tests run, AI reviews diff,\nresults written back to story.\nReady for human review!", ACCENT_GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        col = i % 3
        row = i // 3
        left = Inches(0.6) + Inches(col * 4.1)
        top = Inches(1.5) + Inches(row * 2.8)
        w = Inches(3.7)
        h = Inches(2.3)

        box = add_shape_box(slide, left, top, w, h, WHITE, color)

        # Step number
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, left + Inches(0.8), top + Inches(0.2), w - Inches(1), Inches(0.35),
                     title, font_size=16, bold=True, color=color)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.8), w - Inches(0.4), Inches(1.3),
                     desc, font_size=12, color=DARK_TEXT)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_header(slide, "Roadmap & Next Steps", "Path from prototype to production")

    phases = [
        ("Phase 1 — Current", "Prototype (Completed)", [
            "Plan-Approve-Execute model with per-file review",
            "Smart merge: CLI reads files / API appends safely",
            "Story queue with workspace reset between stories",
            "Dashboard with plan approval modal + SSE",
            "State transitions + write-back to DevOps/Zendesk",
            "Data consent, resume support, alert toasts",
        ], ACCENT_GREEN, "✓"),
        ("Phase 2 — Short Term", "Stabilization (Q2 2026)", [
            "Concurrent execution lock + rebase conflict recovery",
            "Plan approval timeout + empty workspace validation",
            "Dashboard auth + webhook HMAC verification",
            "Comprehensive test suite + CI/CD pipeline",
            "Docker containerization",
        ], ACCENT_BLUE, "→"),
        ("Phase 3 — Medium Term", "Scale (Q3-Q4 2026)", [
            "Multi-team & multi-project support",
            "PR auto-creation and management",
            "Learning from past implementations",
            "Metrics & analytics dashboard",
            "Slack / Teams integration",
        ], ACCENT_PURPLE, "◎"),
    ]

    for i, (phase_title, phase_sub, items, color, icon) in enumerate(phases):
        left = Inches(0.5) + Inches(i * 4.15)
        top = Inches(1.5)
        w = Inches(3.85)

        # Phase header
        header = add_shape_box(slide, left, top, w, Inches(0.9), color, None)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.08), w - Inches(0.3), Inches(0.35),
                     f"{icon}  {phase_title}", font_size=15, bold=True, color=WHITE)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.45), w - Inches(0.3), Inches(0.3),
                     phase_sub, font_size=11, color=RGBColor(0xFF, 0xFF, 0xFF))

        # Items
        item_box = add_shape_box(slide, left, Inches(2.45), w, Inches(4.2), WHITE, color)
        for j, item in enumerate(items):
            y = Inches(2.6) + Inches(j * 0.55)
            add_text_box(slide, left + Inches(0.2), y, w - Inches(0.4), Inches(0.45),
                         f"•  {item}", font_size=11, color=DARK_TEXT)

    # Ask box
    ask_box = add_shape_box(slide, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.7),
                            LIGHT_BLUE_BG, ACCENT_BLUE)
    add_text_box(slide, Inches(0.8), Inches(6.88), Inches(11.73), Inches(0.5),
                 "Ask:  Approval to dedicate 2 sprints for Phase 2 stabilization and team onboarding.",
                 font_size=15, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Accent band at top
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT_BLUE
    band.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.8),
                 "Thank You", font_size=40, bold=True, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.6),
                 "DevOps AI Agent — Automating the Tedious, Empowering the Developer",
                 font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.7), Inches(2.33), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = BORDER_GRAY
    div.line.fill.background()

    # Summary points
    summary = [
        "Plan-Approve-Execute: Safe, human-in-the-loop AI implementation",
        "Smart merge: CLI reads files directly, API appends without wiping code",
        "Dashboard with plan approval, story queue, and real-time progress",
        "Ready for Phase 2 hardening and team onboarding",
    ]
    for i, item in enumerate(summary):
        y = Inches(4.0) + Inches(i * 0.45)
        add_text_box(slide, Inches(3.5), y, Inches(6.5), Inches(0.35),
                     f"✓  {item}", font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.4),
                 "Questions & Discussion", font_size=20, bold=True, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Contact
    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.35),
                 "Pankaj Raundal  |  VTD Producer Connectors  |  devops-ai-agent",
                 font_size=12, color=SUBTITLE_TEXT, alignment=PP_ALIGN.CENTER)

    # Bottom accent band
    band2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12))
    band2.fill.solid()
    band2.fill.fore_color.rgb = ACCENT_BLUE
    band2.line.fill.background()


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_features(prs)
    slide_dashboard(prs)
    slide_benefits(prs)
    slide_demo(prs)
    slide_roadmap(prs)
    slide_closing(prs)

    out_path = "DevOps_AI_Agent_Presentation.pptx"
    prs.save(out_path)
    print(f"✓ Presentation saved to: {out_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
